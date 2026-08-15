"""
Builds the Krypsis class-presentation deck (Krypsis_Presentation.pptx) from
the same 6-slide structure covered in the syllabus-mapped talk: Introduction,
Application, Tools Used, Methodology, Status & Results, Learning.

Run: ..\\venv\\Scripts\\python.exe build_deck.py
Output: presentation/Krypsis_Presentation.pptx
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
OUT_PATH = os.path.join(HERE, "Krypsis_Presentation.pptx")

# Palette, matching the project's web presentation deck.
INK = RGBColor(0x17, 0x24, 0x20)
INK_SOFT = RGBColor(0x56, 0x65, 0x5D)
PAPER = RGBColor(0xF2, 0xF4, 0xF0)
PAPER_DEEP = RGBColor(0xE7, 0xEB, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x14, 0x6B, 0x62)
ACCENT_SOFT = RGBColor(0xD7, 0xEB, 0xE7)
WARM = RGBColor(0xC9, 0x7A, 0x2B)
WARM_SOFT = RGBColor(0xF5, 0xE2, 0xC8)
GOOD = RGBColor(0x2E, 0x7D, 0x4F)
LINE = RGBColor(0xCF, 0xD6, 0xCB)

FONT_DISPLAY = "Arial Black"
FONT_BODY = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs, bg=PAPER):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()
    bg_shape.shadow.inherit = False
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(bg_shape._element)
    spTree.insert(2, bg_shape._element)
    return slide


def add_text(slide, left, top, width, height, text, size=18, color=INK,
             bold=False, font=FONT_BODY, align=PP_ALIGN.LEFT, italic=False,
             line_spacing=1.15, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=15, color=INK,
                 font=FONT_BODY, bullet_color=ACCENT, space_after=10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.2
        run = p.add_run()
        run.text = "▪  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_tag(slide, left, top, text, bg=ACCENT_SOFT, fg=ACCENT):
    w, h = Inches(2.6), Inches(0.42)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = FONT_BODY
    return shape


def add_card(slide, left, top, width, height, bg=WHITE, line=LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.05
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_stat_chip(slide, left, top, width, height, number, label, color=ACCENT):
    add_card(slide, left, top, width, height)
    add_text(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.55),
              number, size=28, color=color, bold=True, font=FONT_DISPLAY)
    add_text(slide, left + Inches(0.15), top + height - Inches(0.55), width - Inches(0.3), Inches(0.5),
              label, size=11, color=INK_SOFT)


def title_slide(prs):
    slide = blank_slide(prs, bg=PAPER_DEEP)
    add_text(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(0.4),
              "COMPUTER NETWORKS — COURSE PROJECT PRESENTATION",
              size=13, color=WARM, bold=True)
    add_text(slide, Inches(0.75), Inches(2.0), Inches(11), Inches(2),
              "KRYPSIS", size=96, color=INK, bold=True, font=FONT_DISPLAY)
    add_text(slide, Inches(0.8), Inches(3.85), Inches(10.5), Inches(0.9),
              "Federated Learning–based Network Intrusion Detection with a\nCustom Communication Protocol",
              size=20, color=INK_SOFT)
    add_text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.9),
              "Rithika K · CB.SC.U4AIE25126      Pradhanya S · CB.SC.U4AIE25148\n"
              "Sathya K · CB.SC.U4AIE25154      Vinudharshini PP · CB.SC.U4AIE25161",
              size=12, color=INK_SOFT)


def slide_intro(prs):
    slide = blank_slide(prs)
    add_tag(slide, Inches(0.7), Inches(0.55), "01  INTRODUCTION")
    add_text(slide, Inches(0.7), Inches(1.15), Inches(11.5), Inches(0.9),
              "Why this project, and what it is", size=38, bold=True, font=FONT_DISPLAY)

    add_text(slide, Inches(0.7), Inches(2.15), Inches(5.6), Inches(0.4),
              "WHY WE CHOSE IT", size=14, bold=True, color=ACCENT)
    add_bullets(slide, Inches(0.7), Inches(2.6), Inches(5.6), Inches(4),
        [
            "Cyberattacks are growing faster than networks can be secured, and most detection systems assume traffic can be pooled onto one central server.",
            "That pooling is often unsafe or illegal — traffic reveals internal structure, and privacy regulations restrict sharing it across organizations.",
            "We wanted a project connecting course networking concepts — protocols, layers, distributed communication — to a live research problem in AI and security.",
        ])

    add_text(slide, Inches(6.7), Inches(2.15), Inches(5.9), Inches(0.4),
              "WHAT THE PROJECT IS", size=14, bold=True, color=ACCENT)
    add_bullets(slide, Inches(6.7), Inches(2.6), Inches(5.9), Inches(4),
        [
            "A Network Intrusion Detection System (NIDS) — watches network connections and classifies each as normal or attack.",
            "Trained using Federated Learning: multiple simulated clients train one shared detector without ever sending raw traffic anywhere.",
            "Extended with a custom communication protocol (in progress) for exchanging model updates — instead of plain HTTP/gRPC.",
        ])


def slide_application(prs):
    slide = blank_slide(prs)
    add_tag(slide, Inches(0.7), Inches(0.55), "02  APPLICATION")
    add_text(slide, Inches(0.7), Inches(1.15), Inches(11.5), Inches(0.9),
              "Where this matters, and what we did", size=36, bold=True, font=FONT_DISPLAY)

    add_text(slide, Inches(0.7), Inches(2.15), Inches(5.6), Inches(0.4),
              "REAL-WORLD APPLICATION", size=14, bold=True, color=ACCENT)
    add_bullets(slide, Inches(0.7), Inches(2.6), Inches(5.6), Inches(4),
        [
            "Hospitals, banks, universities jointly training a shared intrusion detector without exposing internal network layouts to each other.",
            "IoT / edge networks — routers and smart devices each contributing local learning (the “fog computing” pattern from Unit 3).",
            "Cross-border threat intelligence sharing between organizations bound by different data-protection laws.",
        ])

    add_text(slide, Inches(6.7), Inches(2.15), Inches(5.9), Inches(0.4),
              "WHAT WE ACTUALLY BUILT", size=14, bold=True, color=ACCENT)
    add_bullets(slide, Inches(6.7), Inches(2.6), Inches(5.9), Inches(4),
        [
            "Simulated 5 independent clients, each holding its own slice of network traffic data.",
            "Trained one shared detection model across all 5 using Federated Averaging (FedAvg).",
            "Tested both similar-traffic clients (IID) and very different-traffic clients (non-IID) — mimicking real organizational diversity.",
        ])


def slide_tools(prs):
    slide = blank_slide(prs)
    add_tag(slide, Inches(0.7), Inches(0.55), "03  TOOLS USED")
    add_text(slide, Inches(0.7), Inches(1.15), Inches(11.5), Inches(0.9),
              "The toolchain behind it", size=38, bold=True, font=FONT_DISPLAY)
    add_text(slide, Inches(0.7), Inches(1.95), Inches(11.7), Inches(0.6),
              "A machine-learning implementation project, not a network-topology simulation — "
              "so the tools are a data-science stack, not Mininet/Packet Tracer. A deliberate distinction, not a gap.",
              size=13, color=INK_SOFT, italic=True)

    tools = [
        ("Python", "Core implementation language"),
        ("TensorFlow / Keras", "Building & training the neural network"),
        ("scikit-learn", "Preprocessing, metrics, evaluation"),
        ("pandas / NumPy", "Data handling & arrays"),
        ("NSL-KDD dataset", "Benchmark network-traffic records"),
        ("Git & GitHub", "Version control, full project history"),
    ]
    cols, col_w, gap = 3, Inches(3.75), Inches(0.25)
    start_x, start_y, card_h = Inches(0.7), Inches(2.75), Inches(1.15)
    for i, (name, role) in enumerate(tools):
        r, c = divmod(i, cols)
        x = start_x + c * (col_w + gap)
        y = start_y + r * (card_h + gap)
        add_card(slide, x, y, col_w, card_h, bg=WHITE)
        add_text(slide, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.4), Inches(0.4),
                  name, size=16, bold=True)
        add_text(slide, x + Inches(0.2), y + Inches(0.62), col_w - Inches(0.4), Inches(0.45),
                  role, size=11, color=INK_SOFT)

    note_y = Inches(5.35)
    add_card(slide, Inches(0.7), note_y, Inches(11.9), Inches(1.4), bg=WARM_SOFT, line=WARM)
    add_text(slide, Inches(1.0), note_y + Inches(0.15), Inches(0.5), Inches(0.5), "01", size=12, color=WARM, bold=True)
    add_text(slide, Inches(1.0), note_y + Inches(0.15), Inches(11.3), Inches(0.35),
              "FORWARD LINK TO UNIT 3", size=12, bold=True, color=WARM)
    add_text(slide, Inches(1.0), note_y + Inches(0.55), Inches(11.3), Inches(0.75),
              "A natural next step is running the custom protocol over an actual simulated network topology in "
              "Mininet, using OpenFlow to control switch behavior — directly extending this project into the SDN material from Unit 3.",
              size=12, color=INK)


def slide_methodology(prs):
    slide = blank_slide(prs)
    add_tag(slide, Inches(0.7), Inches(0.55), "04  METHODOLOGY")
    add_text(slide, Inches(0.7), Inches(1.15), Inches(11.5), Inches(0.9),
              "How it was built, step by step", size=36, bold=True, font=FONT_DISPLAY)

    steps = [
        "Preprocess NSL-KDD: encode protocol/service/flag, scale numeric features",
        "Split data across 5 simulated clients (IID and non-IID)",
        "Define the model: MLP, 256→128→64→1 neurons",
        "Run FedAvg: local training → weighted averaging → repeat",
    ]
    col_w, gap = Inches(2.85), Inches(0.15)
    x0, y0, h = Inches(0.7), Inches(2.1), Inches(1.3)
    for i, s in enumerate(steps):
        x = x0 + i * (col_w + gap)
        add_card(slide, x, y0, col_w, h, bg=WHITE)
        add_text(slide, x + Inches(0.15), y0 + Inches(0.1), col_w - Inches(0.3), Inches(0.3),
                  f"STEP {i+1}", size=10, bold=True, color=ACCENT)
        add_text(slide, x + Inches(0.15), y0 + Inches(0.4), col_w - Inches(0.3), Inches(0.85),
                  s, size=11.5, color=INK)

    add_text(slide, Inches(0.7), Inches(3.75), Inches(6), Inches(0.4),
              "THE FEDAVG LOOP, IN WORDS", size=14, bold=True, color=ACCENT)
    add_card(slide, Inches(0.7), Inches(4.2), Inches(11.9), Inches(2.7), bg=WHITE)
    add_bullets(slide, Inches(1.0), Inches(4.4), Inches(11.3), Inches(2.4),
        [
            "Server sends the current shared model to all 5 clients.",
            "Each client trains it further on its own data only — nothing else moves.",
            "Clients send back only the updated numbers (weights), never raw traffic.",
            "Server combines everyone's weights into one improved model, weighted by how much data each client had.",
            "Repeat for 15 rounds, checking accuracy after every round.",
        ], size=14)


def slide_status(prs):
    slide = blank_slide(prs)
    add_tag(slide, Inches(0.7), Inches(0.55), "05  STATUS & RESULTS")
    add_text(slide, Inches(0.7), Inches(1.15), Inches(11.5), Inches(0.9),
              "What's done, and the real numbers", size=34, bold=True, font=FONT_DISPLAY)

    stats = [
        ("83.7%", "Centralized baseline accuracy", GOOD),
        ("80.4%", "Federated accuracy, IID clients", ACCENT),
        ("79.0%", "Federated accuracy, non-IID clients", ACCENT),
        ("99.0%", "In-distribution diagnostic", GOOD),
    ]
    col_w, gap = Inches(2.85), Inches(0.15)
    x0, y0, h = Inches(0.7), Inches(2.05), Inches(1.3)
    for i, (num, lbl, col) in enumerate(stats):
        x = x0 + i * (col_w + gap)
        add_stat_chip(slide, x, y0, col_w, h, num, lbl, color=col)

    rows = [
        ("Phase 1-2 — Environment & dataset setup", "Complete", GOOD),
        ("Phase 3 — Preprocessing", "Complete", GOOD),
        ("Phase 4 — Client simulation (IID + non-IID)", "Complete", GOOD),
        ("Phase 5 — Model design & tuning", "Complete", GOOD),
        ("Phase 6 — Federated training loop (FedAvg)", "Complete", GOOD),
        ("Phase 7 — Custom security-fused protocol", "In progress", WARM),
    ]
    ty = Inches(3.7)
    add_card(slide, Inches(0.7), ty, Inches(11.9), Inches(3.3), bg=WHITE)
    row_h = Inches(0.5)
    for i, (name, status, col) in enumerate(rows):
        ry = ty + Inches(0.25) + i * row_h
        add_text(slide, Inches(1.0), ry, Inches(8.5), row_h, name, size=13.5)
        add_text(slide, Inches(9.7), ry, Inches(2.7), row_h, status, size=13.5, bold=True, color=col)


def slide_learning(prs):
    slide = blank_slide(prs)
    add_tag(slide, Inches(0.7), Inches(0.55), "06  LEARNING")
    add_text(slide, Inches(0.7), Inches(1.15), Inches(11.5), Inches(0.9),
              "How this connects to our syllabus", size=34, bold=True, font=FONT_DISPLAY)
    add_text(slide, Inches(0.7), Inches(1.95), Inches(11.7), Inches(0.5),
              "Every unit of this course shows up somewhere in this project — intrusion detection is fundamentally a networking problem.",
              size=13, color=INK_SOFT, italic=True)

    cards = [
        ("UNIT 1 — NETWORK EDGE, OSI LAYERS, PACKET SWITCHING",
         "Every NSL-KDD record is packet/flow-level data: protocol_type (Network layer), "
         "flag (Transport layer), service (Application layer). Our federated clients act as network edge nodes."),
        ("UNIT 2 — APPLICATION LAYER, TRANSPORT LAYER, SOCKETS, ROUTING",
         "Our custom protocol design is application-layer protocol design (vs. HTTP/gRPC). "
         "Client-server exchange mirrors socket-based communication. dst_host_count ties to routing/host addressing."),
        ("UNIT 3 — IOT, FOG/CLOUD COMPUTING, SDN",
         "Federated Learning's architecture IS fog computing: local processing at the edge, summarized results to a "
         "central aggregator. Our custom protocol mirrors SDN/OpenFlow — customizing behavior beyond fixed defaults."),
    ]
    y0, h, gap = Inches(2.65), Inches(1.4), Inches(0.2)
    for i, (unit, body) in enumerate(cards):
        y = y0 + i * (h + gap)
        add_card(slide, Inches(0.7), y, Inches(11.9), h, bg=WHITE)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(0.08), h)
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
        add_text(slide, Inches(1.0), y + Inches(0.15), Inches(11.3), Inches(0.3), unit, size=12, bold=True, color=ACCENT)
        add_text(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.85), body, size=12.5, color=INK)


def slide_roadmap(prs):
    slide = blank_slide(prs)
    add_tag(slide, Inches(0.7), Inches(0.55), "07  ROADMAP")
    add_text(slide, Inches(0.7), Inches(1.15), Inches(11.5), Inches(0.9),
              "What's left, and exactly how", size=36, bold=True, font=FONT_DISPLAY)
    add_text(slide, Inches(0.7), Inches(1.95), Inches(11.7), Inches(0.5),
              "The first half (Objective 1) is fully complete and proven. The remaining work is scoped precisely, not vaguely.",
              size=13, color=INK_SOFT, italic=True)

    items = [
        ("Integrity + fingerprint protocol",
         "Attach a lightweight integrity tag and a compact statistical fingerprint "
         "(per-layer L2 norm + cosine similarity to consensus) to every client update."),
        ("One poisoning attack",
         "Simulate label-flipping attacks from malicious clients to have something real "
         "to test detection against."),
        ("Global vs. Mondrian threshold comparison",
         "The core experiment: does one global anomaly threshold wrongly flag honest "
         "non-IID clients, and does a per-cluster threshold fix that without losing real detection?"),
    ]
    y0, h, gap = Inches(2.6), Inches(1.35), Inches(0.2)
    for i, (title, body) in enumerate(items):
        y = y0 + i * (h + gap)
        add_card(slide, Inches(0.7), y, Inches(11.9), h, bg=WHITE)
        add_text(slide, Inches(1.0), y + Inches(0.15), Inches(0.6), Inches(0.5),
                  str(i + 1), size=22, bold=True, color=ACCENT, font=FONT_DISPLAY)
        add_text(slide, Inches(1.7), y + Inches(0.15), Inches(10.6), Inches(0.3), title, size=14, bold=True)
        add_text(slide, Inches(1.7), y + Inches(0.5), Inches(10.6), Inches(0.75), body, size=12, color=INK_SOFT)


def build():
    prs = new_presentation()
    title_slide(prs)
    slide_intro(prs)
    slide_application(prs)
    slide_tools(prs)
    slide_methodology(prs)
    slide_status(prs)
    slide_learning(prs)
    slide_roadmap(prs)
    prs.save(OUT_PATH)
    print(f"Saved {OUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
